import io
import os

import streamlit as st
from dotenv import load_dotenv
from pptx import Presentation

import gamma_export
from graph import build_graph

load_dotenv()

st.set_page_config(page_title="PitchSnitch", page_icon="🕵️", layout="wide")

st.title("🕵️ PitchSnitch")
st.subheader("From One Word to a McKinsey-Style Strategy Deck")
st.markdown(
    "Type a brand or product name. Get auto-discovered competitors, evidence-graded "
    "research, and an executive-ready strategy deck with a clear recommendation — "
    "in one run."
)

step_columns = st.columns(4)
pipeline_steps = [
    ("🔍", "Discover", "Finds 3 real competitors via live web search"),
    ("📊", "Research", "5 targeted searches per competitor, source-cited"),
    ("🧠", "Analyze", "PM-grade report: positioning, pricing, perception, risk"),
    ("📑", "Strategize", "McKinsey-style deck with a clear recommendation"),
]
for column, (icon, label, description) in zip(step_columns, pipeline_steps):
    with column:
        st.markdown(f"**{icon} {label}**")
        st.caption(description)

st.divider()


def render_full_report(report: dict, company_name: str) -> None:
    competitor_name = report["competitor_name"]

    st.markdown("##### Source Quality")
    if report["sources"]:
        st.dataframe(
            [
                {
                    "Title": s["title"],
                    "Type": s["source_type"],
                    "Favors": s["favors"],
                    "Evidence": s["evidence_quality"],
                    "Bias Risk": s["bias_risk"],
                    "Insight": s["insight"],
                }
                for s in report["sources"]
            ],
            use_container_width=True,
        )
    else:
        st.write("No sources captured.")

    st.markdown("##### Positioning")
    pos = report["positioning"]
    left, right = st.columns(2)
    with left:
        st.markdown(f"**{company_name}**")
        st.write(f"Problem solved: {pos['company_problem_solved']}")
        st.write(f"Promise: {pos['company_promise']}")
        st.write(f"Segment: {pos['company_segment']}")
    with right:
        st.markdown(f"**{competitor_name}**")
        st.write(f"Problem solved: {pos['competitor_problem_solved']}")
        st.write(f"Promise: {pos['competitor_promise']}")
        st.write(f"Segment: {pos['competitor_segment']}")
    st.caption(pos["positioning_difference"])

    st.markdown("##### Feature Comparison")
    if report["feature_comparison"]:
        st.dataframe(
            [
                {
                    "Dimension": f["dimension"],
                    company_name: f["company_value"],
                    competitor_name: f["competitor_value"],
                    "PM Interpretation": f["pm_interpretation"],
                }
                for f in report["feature_comparison"]
            ],
            use_container_width=True,
        )

    st.markdown("##### Customer Perception")
    perception = report["customer_perception"]
    left, right = st.columns(2)
    with left:
        st.markdown(f"**{company_name} praise**")
        for item in perception["company_praise"]:
            st.markdown(f"- {item}")
        st.markdown(f"**{company_name} complaints**")
        for item in perception["company_complaints"]:
            st.markdown(f"- {item}")
    with right:
        st.markdown(f"**{competitor_name} praise**")
        for item in perception["competitor_praise"]:
            st.markdown(f"- {item}")
        st.markdown(f"**{competitor_name} complaints**")
        for item in perception["competitor_complaints"]:
            st.markdown(f"- {item}")

    st.markdown("##### Pricing / Packaging")
    pricing = report["pricing"]
    st.write(pricing["packaging_notes"])
    st.write(pricing["pricing_complaints"])

    st.markdown("##### Strategic Read")
    strategic = report["strategic_read"]
    st.write(f"Competitor's market belief: {strategic['competitor_market_belief']}")
    st.write(f"Expectation being shaped: {strategic['expectation_being_shaped']}")
    st.write(f"{company_name} vulnerability: {strategic['company_vulnerability']}")
    st.write(f"{company_name} advantage: {strategic['company_advantage']}")
    st.write(f"Watch (6-12 months): {strategic['watch_next_6_12_months']}")

    st.markdown("##### PM Recommendations")
    recs = report["recommendations"]
    for label, key in [
        ("Do now", "do_now"),
        ("Do not blindly copy", "do_not_blindly_copy"),
        ("Watch", "watch"),
    ]:
        st.markdown(f"**{label}**")
        for rec in recs[key]:
            st.markdown(f"- {rec['action']} (confidence: {rec['confidence']}) — {rec['rationale']}")

    st.markdown("##### Open Questions")
    for question in report["open_questions"]:
        st.markdown(f"- {question}")


def report_to_markdown(report: dict, company_name: str) -> str:
    competitor_name = report["competitor_name"]
    lines = [f"## {competitor_name}", "", report["executive_summary"], ""]

    lines.append("### Source Quality")
    if report["sources"]:
        lines.append("| Title | Type | Favors | Evidence | Bias Risk | Insight |")
        lines.append("|---|---|---|---|---|---|")
        for s in report["sources"]:
            lines.append(
                f"| {s['title']} | {s['source_type']} | {s['favors']} | "
                f"{s['evidence_quality']} | {s['bias_risk']} | {s['insight']} |"
            )
    else:
        lines.append("No sources captured.")
    lines.append("")

    pos = report["positioning"]
    lines.append("### Positioning")
    lines.append(
        f"- **{company_name}** — Problem solved: {pos['company_problem_solved']}; "
        f"Promise: {pos['company_promise']}; Segment: {pos['company_segment']}"
    )
    lines.append(
        f"- **{competitor_name}** — Problem solved: {pos['competitor_problem_solved']}; "
        f"Promise: {pos['competitor_promise']}; Segment: {pos['competitor_segment']}"
    )
    lines.append(f"*{pos['positioning_difference']}*")
    lines.append("")

    lines.append("### Feature Comparison")
    if report["feature_comparison"]:
        lines.append(f"| Dimension | {company_name} | {competitor_name} | PM Interpretation |")
        lines.append("|---|---|---|---|")
        for f in report["feature_comparison"]:
            lines.append(
                f"| {f['dimension']} | {f['company_value']} | {f['competitor_value']} | "
                f"{f['pm_interpretation']} |"
            )
    lines.append("")

    perception = report["customer_perception"]
    lines.append("### Customer Perception")
    lines.append(f"- **{company_name} praise:** " + "; ".join(perception["company_praise"]))
    lines.append(
        f"- **{company_name} complaints:** " + "; ".join(perception["company_complaints"])
    )
    lines.append(
        f"- **{competitor_name} praise:** " + "; ".join(perception["competitor_praise"])
    )
    lines.append(
        f"- **{competitor_name} complaints:** "
        + "; ".join(perception["competitor_complaints"])
    )
    lines.append("")

    pricing = report["pricing"]
    lines.append("### Pricing / Packaging")
    lines.append(f"- Summary: {pricing['summary']}")
    lines.append(f"- Packaging: {pricing['packaging_notes']}")
    lines.append(f"- Complaints: {pricing['pricing_complaints']}")
    lines.append("")

    strategic = report["strategic_read"]
    lines.append("### Strategic Read")
    lines.append(f"- Competitor's market belief: {strategic['competitor_market_belief']}")
    lines.append(f"- Expectation being shaped: {strategic['expectation_being_shaped']}")
    lines.append(f"- {company_name} vulnerability: {strategic['company_vulnerability']}")
    lines.append(f"- {company_name} advantage: {strategic['company_advantage']}")
    lines.append(f"- Watch (6-12 months): {strategic['watch_next_6_12_months']}")
    lines.append("")

    lines.append("### PM Recommendations")
    recs = report["recommendations"]
    for label, key in [
        ("Do now", "do_now"),
        ("Do not blindly copy", "do_not_blindly_copy"),
        ("Watch", "watch"),
    ]:
        lines.append(f"**{label}**")
        for rec in recs[key]:
            lines.append(f"- {rec['action']} (confidence: {rec['confidence']}) — {rec['rationale']}")
    lines.append("")

    lines.append("### Open Questions")
    for question in report["open_questions"]:
        lines.append(f"- {question}")
    lines.append("")

    return "\n".join(lines)


def render_deck(deck: dict) -> None:
    st.write(deck["narrative_spine"])

    for slide in deck["slides"]:
        with st.expander(f"Slide {slide['slide_number']}: {slide['title']}"):
            st.markdown(f"**Objective:** {slide['objective']}")
            st.markdown(f"**Main message:** {slide['main_message']}")
            st.markdown(f"**Recommended visual:** {slide['recommended_visual']}")
            st.markdown("**Slide content:**")
            st.markdown(slide["slide_content"])
            st.markdown(f"**Evidence to cite:** {slide['evidence_to_cite']}")
            st.caption(f"Speaker notes: {slide['speaker_notes']}")

    st.markdown("##### Final Recommendation")
    roadmap = deck["final_recommendation"]
    for label, key in [
        ("Do now", "do_now"),
        ("Do not blindly copy", "do_not_blindly_copy"),
        ("Watch", "watch"),
    ]:
        st.markdown(f"**{label}**")
        for action in roadmap[key]:
            st.markdown(
                f"- {action['action']} (owner: {action['owner_type']}, "
                f"confidence: {action['confidence']}, evidence: "
                f"{action['evidence_quality']}) — {action['rationale']}"
            )

    st.markdown("##### Evidence Caveats")
    st.write(deck["evidence_caveats"])


def _add_bullet_slide(prs: Presentation, title: str, bullets: list, body_text: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.word_wrap = True
    first = True
    if body_text:
        body.text = body_text
        first = False
    for bullet in bullets:
        if first:
            body.text = bullet
            first = False
        else:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 1
    return slide


def deck_to_pptx(deck: dict, company_name: str) -> bytes:
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = deck["deck_title"]
    title_slide.placeholders[1].text = f"Competitive Strategy Deck — {company_name}"

    _add_bullet_slide(prs, "Narrative Spine", [], body_text=deck["narrative_spine"])

    for slide_data in deck["slides"]:
        bullets = [line.strip("-* ").strip() for line in slide_data["slide_content"].split("\n") if line.strip()]
        bullets.append(f"Evidence: {slide_data['evidence_to_cite']}")
        slide = _add_bullet_slide(
            prs,
            f"{slide_data['slide_number']}. {slide_data['title']}",
            bullets,
            body_text=slide_data["main_message"],
        )
        slide.notes_slide.notes_text_frame.text = slide_data["speaker_notes"]

    roadmap = deck["final_recommendation"]
    for label, key in [
        ("Recommendation: Do Now", "do_now"),
        ("Recommendation: Do Not Blindly Copy", "do_not_blindly_copy"),
        ("Recommendation: Watch", "watch"),
    ]:
        bullets = [
            f"{a['action']} (Owner: {a['owner_type']}, Confidence: {a['confidence']}, "
            f"Evidence: {a['evidence_quality']}) — {a['rationale']}"
            for a in roadmap[key]
        ]
        _add_bullet_slide(prs, label, bullets)

    _add_bullet_slide(prs, "Evidence Caveats", [], body_text=deck["evidence_caveats"])

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


company_name = st.text_input(
    "Company / Product Name", placeholder="e.g. Notion, Figma, TurboTax..."
)

if st.button("Run Research Pipeline"):
    if not os.environ.get("OPENAI_API_KEY"):
        st.error("OPENAI_API_KEY is not set. Add it to your .env file.")
        st.stop()
    if not company_name:
        st.warning("Enter a company or product name.")
        st.stop()

    graph = build_graph()

    with st.status(f"Researching competitors for {company_name}...", expanded=True) as status:
        st.write("Discovering top competitors...")
        result = graph.invoke({
            "company": company_name,
            "competitor_queue": [],
            "final_reports": [],
        })
        st.write(f"Compiled {len(result['final_reports'])} competitor reports.")
        status.update(label="Research complete", state="complete")

    st.session_state["result"] = result
    st.session_state["company_name"] = company_name
    st.session_state.pop("gamma_result", None)
    st.session_state.pop("gamma_pptx_bytes", None)
    st.session_state.pop("gamma_error", None)

if "result" in st.session_state:
    result = st.session_state["result"]
    company_name = st.session_state["company_name"]
    reports = result["final_reports"]

    st.markdown("### Executive Summary")
    st.write(result["overall_summary"])

    full_markdown = (
        f"# Competitive Analysis: {company_name}\n\n"
        f"## Executive Summary\n\n{result['overall_summary']}\n\n"
        + "\n".join(report_to_markdown(r, company_name) for r in reports)
    )
    st.download_button(
        label="Download Full Report",
        data=full_markdown,
        file_name=f"{company_name.replace(' ', '_')}_competitive_analysis.md",
        mime="text/markdown",
    )

    columns = st.columns(len(reports))
    for column, report in zip(columns, reports):
        with column:
            with st.container(border=True):
                st.markdown(f"#### 📦 {report['competitor_name']}")
                st.write(report["executive_summary"])
                st.markdown(f"💰 **Pricing:** {report['pricing']['summary']}")
                st.caption(report["positioning"]["positioning_difference"])

                with st.expander("Full report"):
                    render_full_report(report, company_name)

    st.markdown("### Strategy Deck")
    deck_pptx = deck_to_pptx(result["deck"], company_name)
    st.download_button(
        label="Download Strategy Deck (.pptx)",
        data=deck_pptx,
        file_name=f"{company_name.replace(' ', '_')}_strategy_deck.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    st.markdown("##### Designer-Quality Deck (Gamma)")
    st.caption("Uses your Gamma credits, separate from the .pptx export above.")

    if st.button("Generate Designer Deck with Gamma"):
        gamma_key = os.environ.get("GAMMA_API_KEY")
        if not gamma_key:
            st.error("GAMMA_API_KEY is not set. Add it to your .env file.")
            st.stop()

        st.session_state.pop("gamma_result", None)
        st.session_state.pop("gamma_pptx_bytes", None)
        st.session_state.pop("gamma_error", None)

        try:
            input_text = gamma_export.build_input_text(result["deck"], company_name)
            instructions = gamma_export.build_additional_instructions(result["deck"])
            with st.status("Generating deck with Gamma...", expanded=True) as gamma_status:
                st.write("Submitting content to Gamma...")
                generation_id, warnings = gamma_export.create_generation(
                    input_text, gamma_key, instructions
                )
                if warnings:
                    st.write(f"Gamma warnings: {warnings}")
                st.write(f"Generation started (id: {generation_id}). Polling for completion...")
                progress_line = st.empty()
                gamma_data = gamma_export.poll_generation(
                    generation_id, gamma_key, on_progress=progress_line.write
                )
                progress_line.empty()
                st.write("Downloading export from Gamma...")
                pptx_bytes = gamma_export.download_export(gamma_data["exportUrl"])
                gamma_status.update(label="Gamma deck ready", state="complete")
            st.session_state["gamma_result"] = gamma_data
            st.session_state["gamma_pptx_bytes"] = pptx_bytes
        except gamma_export.GammaAPIError as e:
            st.session_state["gamma_error"] = {"status_code": e.status_code, "message": str(e)}
        except gamma_export.GammaGenerationFailedError as e:
            st.session_state["gamma_error"] = {
                "status_code": None,
                "message": f"Gamma generation failed: {e}",
            }
        except gamma_export.GammaTimeoutError as e:
            st.session_state["gamma_error"] = {
                "status_code": None,
                "message": f"Timed out waiting for Gamma: {e}",
            }

    if "gamma_error" in st.session_state:
        err = st.session_state["gamma_error"]
        code = err["status_code"]
        if code == 401:
            st.error(
                "Gamma rejected the API key (401). Check GAMMA_API_KEY and that the "
                "workspace plan supports the API."
            )
        elif code == 429:
            st.warning("Gamma rate limit hit (429). Wait a bit and try again.")
        elif code is None:
            st.error(f"Could not complete the Gamma request: {err['message']}")
        else:
            st.error(f"Gamma API error ({code}): {err['message']}")

    if "gamma_result" in st.session_state:
        gamma_data = st.session_state["gamma_result"]
        st.success("Designer deck ready.")
        st.link_button("Open in Gamma", gamma_data["gammaUrl"])
        st.download_button(
            label="Download Designer Deck (.pptx)",
            data=st.session_state["gamma_pptx_bytes"],
            file_name=f"{company_name.replace(' ', '_')}_gamma_deck.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        credits = gamma_data.get("credits", {})
        st.caption(
            f"Credits used: {credits.get('deducted', '?')} — "
            f"remaining: {credits.get('remaining', '?')}"
        )

    st.markdown(f"#### {result['deck']['deck_title']}")
    render_deck(result["deck"])
